from __future__ import annotations

import shutil
from pathlib import Path

import pytest
from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Pt

from decklint.cli import main
from decklint.model import load_deck
from decklint.render import render_deck
from decklint.rules import audit_deck


def create_real_pptx(path: Path) -> Path:
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    slide.shapes.title.text = "Real Office Fixture"
    slide.placeholders[1].text = "Native text rendered by LibreOffice"
    presentation.save(path)
    return path


def test_parser_reads_real_python_pptx_package(tmp_path: Path) -> None:
    source = create_real_pptx(tmp_path / "real.pptx")

    deck = load_deck(source)

    assert len(deck.slides) == 1
    assert deck.slides[0].title == "Real Office Fixture"
    assert sum(1 for shape in deck.slides[0].shapes if shape.text) == 2
    assert not deck.broken_relationships


@pytest.mark.skipif(not shutil.which("soffice"), reason="LibreOffice is not installed")
def test_libreoffice_renderer_produces_matching_real_preview(tmp_path: Path) -> None:
    source = create_real_pptx(tmp_path / "real.pptx")
    deck = load_deck(source)

    result = render_deck(deck, source=source, renderer="libreoffice")

    assert result.status == "ok", result.detail
    assert result.used == "libreoffice"
    assert len(result.previews) == len(deck.slides)


def test_default_gate_does_not_fail_only_for_common_last_editor_metadata(tmp_path: Path) -> None:
    source = create_real_pptx(tmp_path / "real.pptx")

    exit_code = main(
        ["audit", str(source), "--renderer", "wireframe", "--output", str(tmp_path / "report")]
    )

    assert exit_code == 0


def test_grouped_text_is_flattened_with_effective_geometry_and_font_size(tmp_path: Path) -> None:
    source = tmp_path / "grouped.pptx"
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    group = slide.shapes.add_group_shape()
    text_box = group.shapes.add_textbox(Inches(2), Inches(1), Inches(4), Inches(1))
    text_box.text = "Grouped six point text"
    text_box.text_frame.paragraphs[0].runs[0].font.size = Pt(6)
    presentation.save(source)

    deck = load_deck(source)
    findings = audit_deck(deck, profile="baseline")

    assert len(deck.slides[0].shapes) == 1
    assert deck.slides[0].shapes[0].bbox.w > 0
    assert deck.slides[0].shapes[0].font_sizes == [6.0]
    assert "readability.small-font" in {finding.rule_id for finding in findings}


def test_grouped_picture_keeps_accessibility_audit(tmp_path: Path) -> None:
    image_path = tmp_path / "pixel.png"
    Image.new("RGB", (32, 32), "red").save(image_path)
    source = tmp_path / "grouped-picture.pptx"
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    group = slide.shapes.add_group_shape()
    picture = group.shapes.add_picture(str(image_path), Inches(1), Inches(1), Inches(2), Inches(2))
    picture._element.nvPicPr.cNvPr.set("descr", "")
    presentation.save(source)

    deck = load_deck(source)
    findings = audit_deck(deck, profile="baseline")

    assert [shape.kind for shape in deck.slides[0].shapes] == ["picture"]
    assert "accessibility.missing-alt-text" in {finding.rule_id for finding in findings}
